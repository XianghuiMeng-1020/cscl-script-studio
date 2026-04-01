"""Document service for course materials management.
S2.3: PDF extraction uses pypdf page text only; binary marker guardrails;
normalize_text (control chars, NFKC, merge whitespace); structured error codes.
M1: PDF extraction runs with configurable timeout; returns EXTRACTION_TIMEOUT on timeout.
"""
import os
import re
import logging
import unicodedata

logger = logging.getLogger(__name__)
import hashlib
import uuid
import io
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError
from typing import Dict, Any, List, Optional, Tuple
from werkzeug.utils import secure_filename
from sqlalchemy.exc import IntegrityError
from app.db import db
from app.models import CSCLCourseDocument, CSCLDocumentChunk
from app.config import Config

# Error codes for upload API
PDF_PARSE_FAILED = 'PDF_PARSE_FAILED'
EMPTY_EXTRACTED_TEXT = 'EMPTY_EXTRACTED_TEXT'
TEXT_TOO_SHORT = 'TEXT_TOO_SHORT'
UNSUPPORTED_FILE_TYPE = 'UNSUPPORTED_FILE_TYPE'
EXTRACTION_TIMEOUT = 'EXTRACTION_TIMEOUT'
EXTRACTION_FAILED = 'EXTRACTION_FAILED'

try:
    from pypdf import PdfReader
    PDF_SUPPORT = True
except ImportError:
    PDF_SUPPORT = False

try:
    import pdfplumber
    PDFPLUMBER_SUPPORT = True
except ImportError:
    PDFPLUMBER_SUPPORT = False

try:
    from pptx import Presentation
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

try:
    from docx import Document as DocxDocument
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

# PDF binary/structure markers (S2.13)
# Only match unambiguous PDF internal structure tokens — avoid false positives on
# normal English words like "stream", "object", "trailer".
_PDF_STRUCTURE_RE = re.compile(
    r'%PDF-\d|^\s*\d+\s+\d+\s+obj\b|\bendobj\b|\bstartxref\b',
    re.IGNORECASE | re.MULTILINE
)
# Legacy marker pattern (kept for normalize_text line filter)
# Only match lines that look like actual PDF structure, not normal English words.
# Require at least 2 markers on the same line, or unambiguous markers like %PDF- / startxref.
_PDF_BINARY_MARKERS = re.compile(
    r'%PDF-\d|\bstartxref\b|^\s*\d+\s+\d+\s+obj\b|\bendobj\b.*\bendobj\b',
    re.MULTILINE | re.IGNORECASE
)
# S2.13: non-printable ratio threshold for "probably binary"
_NON_PRINTABLE_RATIO_THRESHOLD = 0.10
# S2.13: minimum printable ratio per line in sanitize (stricter)
_SANITIZE_PRINTABLE_RATIO_MIN = 0.7
# Minimum printable ratio for a line to be kept (avoid binary noise)
_PRINTABLE_RATIO_MIN = 0.5
# Max consecutive control-char run to allow (then strip)
_CONTROL_CHAR_RUN_MAX = 2
_MIN_EXTRACT_LEN = 80
_PREVIEW_LEN = 5000


def _detected_type_from_mime(mime_type: str) -> str:
    if not mime_type:
        return 'txt'
    m = (mime_type or '').lower()
    if 'pdf' in m:
        return 'pdf'
    if 'presentation' in m or 'vnd.ms-powerpoint' in m or 'pptx' in m:
        return 'pptx'
    if 'wordprocessing' in m or 'msword' in m or 'docx' in m:
        return 'docx'
    if 'csv' in m or 'spreadsheet' in m:
        return 'csv'
    if 'png' in m or 'jpeg' in m or 'jpg' in m:
        return 'image'
    if 'markdown' in m or 'md' in m:
        return 'txt'
    return 'txt'


def _non_printable_ratio(s: str) -> float:
    """Ratio of non-printable (control/surrogate) chars. 0 = all printable."""
    if not s:
        return 0.0
    n = sum(1 for c in s if unicodedata.category(c)[0] == 'C' and c not in '\n\t')
    return n / len(s)


def is_probably_pdf_binary_text(s: str) -> bool:
    """
    Single server-side detector: True if text looks like PDF binary/structure or garbage.
    - Contains "%PDF-"
    - Tokens: obj/endobj/stream/endstream/xref/trailer/startxref
    - High non-printable ratio
    - Replacement-char (U+FFFD) explosion
    """
    if not s or not s.strip():
        return False
    if '%PDF-' in s:
        return True
    if _PDF_STRUCTURE_RE.search(s):
        return True
    if _non_printable_ratio(s) > _NON_PRINTABLE_RATIO_THRESHOLD:
        return True
    run = 0
    for c in s:
        if unicodedata.category(c)[0] == 'C' and c not in '\n\t' or c == '\ufffd':
            run += 1
            if run >= 20:
                return True
        else:
            run = 0
    return False


def safe_preview_or_none(text: str, max_len: int = _PREVIEW_LEN) -> Optional[str]:
    """
    Single sanitizer before ANY response field that may contain extracted text.
    Returns None if text looks like PDF/binary (caller must return 422, never return original).
    Otherwise returns text truncated to max_len. Use for: extracted_text_preview, content_preview,
    raw_text, document_text, any fallback message that could contain extracted content.
    """
    if not text:
        return ''
    if is_probably_pdf_binary_text(text):
        return None
    return text[:max_len] if len(text) > max_len else text


def to_display_safe_preview(text: str, max_len: int = _PREVIEW_LEN) -> Optional[str]:
    """
    Nuclear filter for API responses: keep ONLY letters, numbers, punctuation, spaces.
    Use for every extracted_text_preview sent to the client. Guarantees no binary/stream
    can ever reach the frontend even if DB or pypdf had edge cases.
    """
    if not text or not isinstance(text, str):
        return None
    if '%PDF-' in text or _PDF_STRUCTURE_RE.search(text):
        return None
    # Keep only printable/safe Unicode categories: Letter, Number, Punctuation, Separator + \t\n
    cleaned = ''.join(
        c for c in text
        if unicodedata.category(c)[0] in ('L', 'N', 'P', 'Z') or c in '\t\n'
    )
    cleaned = re.sub(r'[ \t]+', ' ', re.sub(r'\n{3,}', '\n\n', cleaned)).strip()
    if len(cleaned) < 20:
        return None
    return cleaned[:max_len] if len(cleaned) > max_len else cleaned


def sanitize_extracted_text(s: str) -> str:
    """
    S2.13: NFKC, strip control chars (keep \\n\\t), drop lines with PDF markers or low printable ratio,
    merge blank lines, strip. Do not use raw decode() output for PDF.
    """
    if not s:
        return ''
    s = unicodedata.normalize('NFKC', s)
    s = ''.join(c for c in s if unicodedata.category(c)[0] != 'C' or c in '\n\t')
    s = re.sub(r'[ \t]+', ' ', s)
    s = re.sub(r'\n{3,}', '\n\n', s)
    s = re.sub(r'[\u200b-\u200d\ufeff]', '', s)
    lines = s.split('\n')
    cleaned = []
    for ln in lines:
        if _PDF_BINARY_MARKERS.search(ln):
            continue
        if ln.strip() and _line_printable_ratio(ln) < _SANITIZE_PRINTABLE_RATIO_MIN:
            continue
        cleaned.append(ln)
    s = '\n'.join(cleaned)
    s = re.sub(r'\n{3,}', '\n\n', s)
    return s.strip()


def _line_printable_ratio(line: str) -> float:
    if not line:
        return 1.0
    n = sum(1 for c in line if unicodedata.category(c)[0] in ('L', 'N', 'P', 'Z') or c in ' \t\n')
    return n / len(line)


class DocumentService:
    """Service for managing course documents"""
    
    ALLOWED_EXTENSIONS = {'txt', 'md', 'pdf', 'pptx', 'docx', 'csv', 'png', 'jpg', 'jpeg'}

    @property
    def MAX_FILE_SIZE(self) -> int:
        """M1: Configurable max upload size (bytes)."""
        try:
            from flask import current_app
            mb = float(current_app.config.get('DOCUMENT_MAX_FILE_SIZE_MB', 10))
            return int(mb * 1024 * 1024)
        except Exception:
            return 10 * 1024 * 1024  # 10MB default
    
    def __init__(self):
        self.upload_dir = os.path.join(Config.DATA_DIR, 'course_documents')
        os.makedirs(self.upload_dir, exist_ok=True)
    
    @staticmethod
    def _cleanup_temp_file(file_path):
        """Remove temporary file from disk (best-effort, never raises)."""
        if file_path:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except OSError:
                pass

    @staticmethod
    def _store_file_data(doc_id, file_content):
        """Store file binary in the file_data column via raw SQL (column may not exist)."""
        try:
            from sqlalchemy import text
            db.session.execute(
                text("UPDATE cscl_course_documents SET file_data = :data WHERE id = :id"),
                {"data": file_content, "id": doc_id}
            )
            db.session.commit()
        except Exception:
            db.session.rollback()
            logger.warning("Could not store file_data for doc %s (column may not exist)", doc_id)
    
    def allowed_file(self, filename: str) -> bool:
        """Check if file extension is allowed"""
        if '.' not in filename:
            return False
        ext = filename.rsplit('.', 1)[1].lower()
        return ext in self.ALLOWED_EXTENSIONS
    
    def compute_checksum(self, content: bytes) -> str:
        """Compute SHA256 checksum of content"""
        return hashlib.sha256(content).hexdigest()
    
    def chunk_text(self, text: str, chunk_size: int = 500) -> List[str]:
        """Split text into chunks"""
        words = text.split()
        chunks = []
        current_chunk = []
        current_length = 0
        
        for word in words:
            word_length = len(word) + 1  # +1 for space
            if current_length + word_length > chunk_size and current_chunk:
                chunks.append(' '.join(current_chunk))
                current_chunk = [word]
                current_length = word_length
            else:
                current_chunk.append(word)
                current_length += word_length
        
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks
    
    def clean_text(self, text: str) -> str:
        """
        Clean text: UTF-8 normalization, remove control characters, 
        normalize full-width/half-width, preserve valid punctuation
        (Deprecated: use normalize_text instead)
        """
        return self.normalize_text(text)
    
    def _has_pdf_binary_markers(self, text: str) -> bool:
        """True if text contains typical PDF binary/structure markers (extraction failure)."""
        return bool(_PDF_BINARY_MARKERS.search(text))

    def _pdf_max_pages(self) -> int:
        """M1: Configurable max PDF pages to process."""
        try:
            from flask import current_app
            return int(current_app.config.get('PDF_MAX_PAGES', 500))
        except Exception:
            return 500

    def _extract_with_pypdf(self, data: bytes) -> Tuple[Optional[str], int]:
        """Extract text using pypdf. Returns (cleaned_text, page_count)."""
        if not PDF_SUPPORT:
            return None, 0
        reader = PdfReader(io.BytesIO(data))
        max_pages = self._pdf_max_pages()
        parts = []
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                break
            try:
                t = page.extract_text()
                if t:
                    parts.append(t)
            except Exception:
                continue
        if not parts:
            return None, len(reader.pages)
        raw = '\n\n'.join(parts)
        clean = sanitize_extracted_text(raw)
        return clean, len(parts)

    def _extract_with_pdfplumber(self, data: bytes) -> Tuple[Optional[str], int]:
        """Extract text using pdfplumber (often better for complex PDFs). Returns (cleaned_text, page_count)."""
        if not PDFPLUMBER_SUPPORT:
            return None, 0
        try:
            with pdfplumber.open(io.BytesIO(data)) as pdf:
                max_pages = self._pdf_max_pages()
                parts = []
                for i, page in enumerate(pdf.pages):
                    if i >= max_pages:
                        break
                    try:
                        t = page.extract_text()
                        if t:
                            parts.append(t)
                    except Exception:
                        continue
                if not parts:
                    return None, len(pdf.pages)
                raw = '\n\n'.join(parts)
                clean = sanitize_extracted_text(raw)
                return clean, len(parts)
        except Exception:
            return None, 0

    def extract_text_from_pdf_bytes(self, data: bytes) -> Dict[str, Any]:
        """
        S2.13: Extract text from PDF using pypdf, with pdfplumber fallback when pypdf
        extracts too little (common for complex layouts). Never fallback to data.decode().
        M1: Limits to first PDF_MAX_PAGES pages to avoid long request hangs.
        Returns: { ok: True, extracted_text, extracted_text_preview } or
                 { ok: False, code, error } (error is readable message, never raw bytes).
        """
        if not PDF_SUPPORT:
            return {'ok': False, 'code': PDF_PARSE_FAILED, 'error': 'PDF library not available'}
        try:
            # Try pypdf first
            clean, page_count = self._extract_with_pypdf(data)
            # If pypdf yields very little for a multi-page PDF, try pdfplumber
            min_chars_per_page = 50
            if (page_count > 1 and clean and len(clean.strip()) < page_count * min_chars_per_page
                    and PDFPLUMBER_SUPPORT):
                clean_plumber, _ = self._extract_with_pdfplumber(data)
                if clean_plumber and len(clean_plumber.strip()) > len((clean or '').strip()):
                    clean = clean_plumber
            # Also try pdfplumber if pypdf failed entirely
            if (not clean or len(clean.strip()) < _MIN_EXTRACT_LEN) and PDFPLUMBER_SUPPORT:
                clean, _ = self._extract_with_pdfplumber(data)
            if not clean or len(clean.strip()) < _MIN_EXTRACT_LEN:
                return {'ok': False, 'code': TEXT_TOO_SHORT, 'error': 'Extracted text is too short or empty. If this is a scanned or image-only PDF, paste the text manually in Step 1.'}
            if is_probably_pdf_binary_text(clean):
                return {'ok': False, 'code': PDF_PARSE_FAILED, 'error': 'PDF parsing failed: binary or invalid content detected'}
            preview = clean[:_PREVIEW_LEN] if len(clean) > _PREVIEW_LEN else clean
            return {
                'ok': True,
                'extracted_text': clean,
                'extracted_text_preview': preview
            }
        except Exception as e:
            try:
                safe_len = len(data)
                prefix = hashlib.sha256(data[: min(256, len(data))]).hexdigest()[:16] if data else ''
                logging.getLogger(__name__).warning(
                    'PDF parse failed: len=%s hash_prefix=%s msg=%s',
                    safe_len, prefix, getattr(e, 'message', str(e))[:200]
                )
            except Exception:
                pass
            return {'ok': False, 'code': PDF_PARSE_FAILED, 'error': 'PDF parsing failed: unable to extract text. The file may be corrupted, encrypted, or image-only.'}

    def extract_text_from_pdf(self, file_bytes: bytes) -> str:
        """
        Extract text from PDF using pypdf page text only. No raw decode of bytes.
        Fails with PDF_PARSE_FAILED if binary markers appear in extracted text.
        """
        if not PDF_SUPPORT:
            raise ValueError("PDF解析库未安装。请安装pypdf: pip install pypdf")
        try:
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)
            text_parts = []
            for page in reader.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                except Exception:
                    continue
            if not text_parts:
                raise ValueError("EMPTY_EXTRACTED_TEXT")
            full_text = '\n\n'.join(text_parts)
            normalized_text = self.normalize_text(full_text)
            if self._has_pdf_binary_markers(normalized_text):
                raise ValueError("PDF_PARSE_FAILED")
            if not normalized_text or not normalized_text.strip():
                raise ValueError("EMPTY_EXTRACTED_TEXT")
            if len(normalized_text.strip()) < _MIN_EXTRACT_LEN:
                raise ValueError("TEXT_TOO_SHORT")
            return normalized_text
        except ValueError:
            raise
        except Exception as e:
            raise ValueError("PDF_PARSE_FAILED")
    
    def extract_text_from_plain(self, file_bytes: bytes, file_path: str) -> str:
        """
        Extract text from plain text file (TXT/MD) with encoding detection
        
        Args:
            file_bytes: File content as bytes
            file_path: File path for extension detection
            
        Returns:
            Extracted text string
        """
        encodings = ['utf-8', 'utf-8-sig', 'gb18030', 'big5', 'latin-1']
        text = None
        last_error = None
        
        for encoding in encodings:
            try:
                text = file_bytes.decode(encoding)
                break
            except UnicodeDecodeError as e:
                last_error = e
                continue
        
        if text is None:
            raise ValueError(UNSUPPORTED_FILE_TYPE)
        normalized = self.normalize_text(text)
        if not normalized or not normalized.strip():
            raise ValueError(EMPTY_EXTRACTED_TEXT)
        if len(normalized.strip()) < _MIN_EXTRACT_LEN:
            raise ValueError(TEXT_TOO_SHORT)
        return normalized
    
    def extract_text_from_pptx(self, data: bytes) -> str:
        """Extract text from PowerPoint (PPTX) using python-pptx."""
        if not PPTX_SUPPORT:
            raise ValueError(UNSUPPORTED_FILE_TYPE)
        try:
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        parts.append(shape.text)
            if not parts:
                raise ValueError(EMPTY_EXTRACTED_TEXT)
            raw = '\n\n'.join(parts)
            normalized = self.normalize_text(raw)
            if not normalized or len(normalized.strip()) < _MIN_EXTRACT_LEN:
                raise ValueError(TEXT_TOO_SHORT)
            return normalized
        except ValueError:
            raise
        except Exception as e:
            logging.getLogger(__name__).warning('PPTX extract failed: %s', str(e)[:200])
            raise ValueError(UNSUPPORTED_FILE_TYPE)
    
    def extract_text_from_docx(self, data: bytes) -> str:
        """Extract text from Word (DOCX) using python-docx."""
        if not DOCX_SUPPORT:
            raise ValueError(UNSUPPORTED_FILE_TYPE)
        try:
            doc = DocxDocument(io.BytesIO(data))
            parts = [p.text for p in doc.paragraphs if p.text]
            if not parts:
                raise ValueError(EMPTY_EXTRACTED_TEXT)
            raw = '\n\n'.join(parts)
            normalized = self.normalize_text(raw)
            if not normalized or len(normalized.strip()) < _MIN_EXTRACT_LEN:
                raise ValueError(TEXT_TOO_SHORT)
            return normalized
        except ValueError:
            raise
        except Exception as e:
            logging.getLogger(__name__).warning('DOCX extract failed: %s', str(e)[:200])
            raise ValueError(UNSUPPORTED_FILE_TYPE)
    
    def _printable_ratio(self, line: str) -> float:
        """Ratio of printable (letter, number, punctuation) chars in line."""
        if not line:
            return 1.0
        printable = sum(1 for c in line if unicodedata.category(c)[0] in ('L', 'N', 'P', 'Z') or c in ' \t\n')
        return printable / len(line)

    def normalize_text(self, text: str) -> str:
        """
        Normalize extracted text: NFKC, remove control chars (keep \\n \\t), merge whitespace,
        remove lines that contain obvious PDF binary noise (%PDF-, xref, trailer, obj, stream, etc.),
        drop lines with low printable ratio, and strip excessive control-char runs.
        Caller must still check _has_pdf_binary_markers and fail extraction if present.
        """
        if not text:
            return ''
        text = unicodedata.normalize('NFKC', text)
        # Strip all control chars except \\n and \\t (avoid garbled PDF text)
        text = ''.join(c for c in text if unicodedata.category(c)[0] != 'C' or c in '\n\t')
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r'[\u200b-\u200d\ufeff]', '', text)
        # Remove lines containing PDF binary/structure noise
        lines = text.split('\n')
        cleaned_lines = [
            ln for ln in lines
            if not _PDF_BINARY_MARKERS.search(ln) and self._printable_ratio(ln) >= _PRINTABLE_RATIO_MIN
        ]
        text = '\n'.join(cleaned_lines)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()
    
    def extract_text_from_file(self, file_path: str, mime_type: str) -> str:
        """
        Extract text from file based on mime type
        Supports: txt, md, pdf
        Returns cleaned UTF-8 text
        
        Args:
            file_path: Path to file
            mime_type: MIME type of file
            
        Returns:
            Extracted and normalized text
            
        Raises:
            ValueError: If file type unsupported or extraction fails
        """
        try:
            # Read file as bytes first
            with open(file_path, 'rb') as f:
                file_bytes = f.read()
            
            if mime_type == 'text/plain' or file_path.endswith('.txt'):
                return self.extract_text_from_plain(file_bytes, file_path)
                
            elif mime_type == 'text/markdown' or file_path.endswith('.md'):
                return self.extract_text_from_plain(file_bytes, file_path)
                
            elif mime_type == 'application/pdf' or file_path.endswith('.pdf'):
                return self.extract_text_from_pdf(file_bytes)
                
            elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document' or file_path.endswith('.docx'):
                raise ValueError(UNSUPPORTED_FILE_TYPE)
            else:
                raise ValueError(UNSUPPORTED_FILE_TYPE)
        except ValueError:
            raise
        except Exception:
            raise ValueError(PDF_PARSE_FAILED if (mime_type or '').lower().startswith('application/pdf') else UNSUPPORTED_FILE_TYPE)
    
    def upload_document(self, course_id: str, title: str, file_content: bytes,
                       filename: str, mime_type: str, uploaded_by: str,
                       material_level: str = 'course', extract_text: bool = True,
                       folder_id: str = None) -> Dict[str, Any]:
        """
        Upload and process a document.
        material_level: 'course' | 'lesson'
        extract_text: if False, only store file and metadata; no text extraction or chunks (for RAG storage only).
        folder_id: optional folder/activity to associate this document with
        Returns:
            {
                'document': {...},
                'chunks_count': int,
                'error': str or None
            }
        """
        if material_level not in ('course', 'lesson'):
            material_level = 'course'
        # Validate file type
        if not self.allowed_file(filename):
            return {
                'document': None,
                'chunks_count': 0,
                'error': f'File type not supported. Allowed: {", ".join(self.ALLOWED_EXTENSIONS)}',
                'error_code': UNSUPPORTED_FILE_TYPE
            }
        
        # Check file size
        if len(file_content) > self.MAX_FILE_SIZE:
            return {
                'document': None,
                'chunks_count': 0,
                'error': f'File too large. Maximum size: {self.MAX_FILE_SIZE / 1024 / 1024}MB',
                'error_code': 'FILE_TOO_LARGE'
            }
        
        # Compute checksum
        checksum = self.compute_checksum(file_content)
        
        # Check for duplicate
        existing = CSCLCourseDocument.query.filter_by(
            course_id=course_id,
            checksum=checksum
        ).first()
        
        if existing:
            doc_dict = existing.to_dict()
            # Only use chunks that are not PDF binary (e.g. legacy bad data)
            safe_chunks = [c.chunk_text for c in existing.chunks if c.chunk_text and not is_probably_pdf_binary_text(c.chunk_text)]
            full_text = ' '.join(safe_chunks) if safe_chunks else ''
            preview = to_display_safe_preview(full_text)
            if preview is None:
                return {
                    'document': None,
                    'chunks_count': 0,
                    'error': 'Unable to parse readable text from PDF. Please upload a text-based PDF or paste plain text.',
                    'error_code': PDF_PARSE_FAILED
                }
            return {
                'document': doc_dict,
                'chunks_count': len(existing.chunks),
                'error': None,
                'extraction_metadata': {
                    'detected_type': _detected_type_from_mime(mime_type),
                    'extracted_char_count': len(full_text),
                    'extraction_method': 'cached',
                    'warnings': [],
                    'extracted_text_preview': preview
                }
            }
        
        safe_filename = secure_filename(filename)
        file_id = str(uuid.uuid4())
        
        # Also write to disk as fallback for text extraction functions that need a file path
        file_path = os.path.join(self.upload_dir, f"{file_id}_{safe_filename}")
        try:
            with open(file_path, 'wb') as f:
                f.write(file_content)
        except OSError:
            file_path = None
        
        if not extract_text:
            document = CSCLCourseDocument(
                id=file_id,
                course_id=course_id,
                folder_id=folder_id,
                title=title or safe_filename,
                source_type='file',
                storage_uri=file_path,
                mime_type=mime_type,
                checksum=checksum,
                file_size=len(file_content),
                material_level=material_level,
                uploaded_by=uploaded_by
            )
            db.session.add(document)
            try:
                db.session.commit()
                self._store_file_data(document.id, file_content)
            except IntegrityError:
                db.session.rollback()
                existing = CSCLCourseDocument.query.filter_by(course_id=course_id, checksum=checksum).first()
                if existing:
                    return {
                        'document': existing.to_dict(),
                        'chunks_count': 0,
                        'error': None,
                        'extraction_metadata': {
                            'detected_type': _detected_type_from_mime(mime_type),
                            'extracted_char_count': 0,
                            'extraction_method': 'skipped',
                            'warnings': [],
                            'extracted_text_preview': None
                        }
                    }
                raise
            self._cleanup_temp_file(file_path)
            return {
                'document': document.to_dict(),
                'chunks_count': 0,
                'error': None,
                'extraction_metadata': {
                    'detected_type': _detected_type_from_mime(mime_type),
                    'extracted_char_count': 0,
                    'extraction_method': 'skipped',
                    'warnings': [],
                    'extracted_text_preview': None
                }
            }
        
        is_pdf = (mime_type or '').lower().startswith('application/pdf') or (filename or '').lower().endswith('.pdf')
        is_pptx = (mime_type or '').lower().count('presentation') or (filename or '').lower().endswith('.pptx')
        is_docx = (mime_type or '').lower().count('wordprocessing') or (mime_type or '').lower().count('msword') or (filename or '').lower().endswith('.docx')
        is_image = (filename or '').lower().endswith(('.png', '.jpg', '.jpeg'))
        text = None
        if is_pdf:
            timeout_sec = getattr(Config, 'PDF_EXTRACTION_TIMEOUT_SECONDS', 60)
            with ThreadPoolExecutor(max_workers=1) as ex:
                fut = ex.submit(self.extract_text_from_pdf_bytes, file_content)
                try:
                    pdf_result = fut.result(timeout=timeout_sec)
                except FuturesTimeoutError:
                    self._cleanup_temp_file(file_path)
                    return {
                        'document': None, 'chunks_count': 0,
                        'error': f'PDF extraction timed out after {timeout_sec}s. Try a smaller file or fewer pages.',
                        'error_code': EXTRACTION_TIMEOUT
                    }
            if not pdf_result.get('ok'):
                self._cleanup_temp_file(file_path)
                return {
                    'document': None, 'chunks_count': 0,
                    'error': pdf_result.get('error', 'PDF extraction failed'),
                    'error_code': pdf_result.get('code', PDF_PARSE_FAILED)
                }
            text = pdf_result['extracted_text']
        elif is_pptx:
            try:
                text = self.extract_text_from_pptx(file_content)
            except ValueError as e:
                self._cleanup_temp_file(file_path)
                err = str(e).strip()
                error_code = EXTRACTION_FAILED if err == EMPTY_EXTRACTED_TEXT or err == TEXT_TOO_SHORT else UNSUPPORTED_FILE_TYPE
                return {'document': None, 'chunks_count': 0, 'error': err, 'error_code': error_code}
            except Exception:
                self._cleanup_temp_file(file_path)
                return {'document': None, 'chunks_count': 0, 'error': 'PPTX extraction failed.', 'error_code': UNSUPPORTED_FILE_TYPE}
        elif is_docx:
            try:
                text = self.extract_text_from_docx(file_content)
            except ValueError as e:
                self._cleanup_temp_file(file_path)
                err = str(e).strip()
                error_code = EXTRACTION_FAILED if err == EMPTY_EXTRACTED_TEXT or err == TEXT_TOO_SHORT else UNSUPPORTED_FILE_TYPE
                return {'document': None, 'chunks_count': 0, 'error': err, 'error_code': error_code}
            except Exception:
                self._cleanup_temp_file(file_path)
                return {'document': None, 'chunks_count': 0, 'error': 'DOCX extraction failed.', 'error_code': UNSUPPORTED_FILE_TYPE}
        elif is_image:
            text = '[Image: ' + (title or safe_filename) + ']'
        else:
            try:
                text = self.extract_text_from_file(file_path, mime_type)
            except ValueError as e:
                self._cleanup_temp_file(file_path)
                err = str(e).strip()
                if err == PDF_PARSE_FAILED: error_code = PDF_PARSE_FAILED
                elif err == EMPTY_EXTRACTED_TEXT: error_code = EMPTY_EXTRACTED_TEXT
                elif err == TEXT_TOO_SHORT: error_code = TEXT_TOO_SHORT
                elif 'DOCX' in err or 'unsupported' in err.lower(): error_code = UNSUPPORTED_FILE_TYPE
                else: error_code = 'EXTRACTION_FAILED'
                return {'document': None, 'chunks_count': 0, 'error': err, 'error_code': error_code}
            except Exception:
                self._cleanup_temp_file(file_path)
                return {
                    'document': None, 'chunks_count': 0,
                    'error': 'Extraction failed',
                    'error_code': PDF_PARSE_FAILED if 'pdf' in (mime_type or '').lower() else 'EXTRACTION_FAILED'
                }
        
        document = CSCLCourseDocument(
            id=file_id,
            course_id=course_id,
            folder_id=folder_id,
            title=title or safe_filename,
            source_type='file',
            storage_uri=file_path,
            mime_type=mime_type,
            checksum=checksum,
            file_size=len(file_content),
            material_level=material_level,
            uploaded_by=uploaded_by
        )

        db.session.add(document)
        db.session.flush()
        
        chunks = self.chunk_text(text)
        chunk_objects = []
        for idx, chunk_text in enumerate(chunks):
            chunk = CSCLDocumentChunk(
                document_id=document.id,
                chunk_index=idx,
                chunk_text=chunk_text,
                token_count=len(chunk_text.split())
            )
            chunk_objects.append(chunk)
            db.session.add(chunk)
        
        full_text = ' '.join(chunks)
        preview = to_display_safe_preview(full_text)
        if preview is None:
            db.session.rollback()
            self._cleanup_temp_file(file_path)
            return {
                'document': None, 'chunks_count': 0,
                'error': 'Unable to parse readable text from PDF. Please upload a text-based PDF or paste plain text.',
                'error_code': PDF_PARSE_FAILED
            }
        
        try:
            db.session.commit()
        except IntegrityError:
            db.session.rollback()
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                except OSError:
                    pass
            existing = CSCLCourseDocument.query.filter_by(
                course_id=course_id,
                checksum=checksum
            ).first()
            if existing:
                doc_dict = existing.to_dict()
                safe_chunks = [c.chunk_text for c in existing.chunks if c.chunk_text and not is_probably_pdf_binary_text(c.chunk_text)]
                full_text_ex = ' '.join(safe_chunks) if safe_chunks else ''
                preview_ex = to_display_safe_preview(full_text_ex)
                if preview_ex is None:
                    preview_ex = ''
                return {
                    'document': doc_dict,
                    'chunks_count': len(existing.chunks),
                    'error': None,
                    'extraction_metadata': {
                        'detected_type': _detected_type_from_mime(mime_type),
                        'extracted_char_count': len(full_text_ex),
                        'extraction_method': 'cached',
                        'warnings': [],
                        'extracted_text_preview': preview_ex
                    }
                }
            raise
        self._store_file_data(document.id, file_content)
        self._cleanup_temp_file(file_path)
        doc_dict = document.to_dict()
        detected = _detected_type_from_mime(mime_type)
        extraction_method = 'pypdf_page_text' if detected == 'pdf' else 'plain_text'
        return {
            'document': doc_dict,
            'chunks_count': len(chunks),
            'error': None,
            'extraction_metadata': {
                'detected_type': detected,
                'extracted_char_count': len(full_text),
                'extraction_method': extraction_method,
                'warnings': [],
                'extracted_text_preview': preview
            }
        }
    
    def upload_text_document(self, course_id: str, title: str, text: str,
                            uploaded_by: str, material_level: str = 'course',
                            folder_id: str = None) -> Dict[str, Any]:
        """
        Upload a text document directly (paste text).
        material_level: 'course' | 'lesson'
        folder_id: optional folder/activity to associate this document with
        Text is cleaned and normalized before processing.
        """
        if material_level not in ('course', 'lesson'):
            material_level = 'course'
        # Clean text first
        cleaned_text = self.normalize_text(text)
        if is_probably_pdf_binary_text(cleaned_text):
            return {
                'document': None,
                'chunks_count': 0,
                'error': 'Pasted content appears to be PDF binary or invalid. Please paste plain text only.',
                'error_code': PDF_PARSE_FAILED
            }
        if not cleaned_text or len(cleaned_text.strip()) < 80:
            return {
                'document': None,
                'chunks_count': 0,
                'error': '文本内容过短（至少80个字符）。请提供有效的课程大纲内容。',
                'error_code': 'TEXT_TOO_SHORT'
            }
        
        # Compute checksum on cleaned text
        checksum = self.compute_checksum(cleaned_text.encode('utf-8'))
        
        # Check for duplicate
        existing = CSCLCourseDocument.query.filter_by(
            course_id=course_id,
            checksum=checksum
        ).first()
        
        if existing:
            return {
                'document': existing.to_dict(),
                'chunks_count': len(existing.chunks),
                'error': None
            }
        
        # Create document record
        document = CSCLCourseDocument(
            course_id=course_id,
            folder_id=folder_id,
            title=title,
            source_type='text',
            storage_uri=None,
            mime_type='text/plain',
            checksum=checksum,
            file_size=len(text.encode('utf-8')),
            material_level=material_level,
            uploaded_by=uploaded_by
        )

        db.session.add(document)
        db.session.flush()

        # Create chunks from cleaned text
        chunks = self.chunk_text(cleaned_text)
        for idx, chunk_text in enumerate(chunks):
            chunk = CSCLDocumentChunk(
                document_id=document.id,
                chunk_index=idx,
                chunk_text=chunk_text,
                token_count=len(chunk_text.split())
            )
            db.session.add(chunk)
        
        db.session.commit()
        doc_dict = document.to_dict()
        full_text = ' '.join(chunks)
        preview = safe_preview_or_none(full_text)
        if preview is None:
            preview = ''  # fallback: never return binary
        return {
            'document': doc_dict,
            'chunks_count': len(chunks),
            'error': None,
            'extraction_metadata': {
                'detected_type': 'txt',
                'extracted_char_count': len(full_text),
                'extraction_method': 'plain_text',
                'warnings': [],
                'extracted_text_preview': preview
            }
        }
    
    def get_course_documents(self, course_id: str, folder_id: str = None) -> List[Dict[str, Any]]:
        """Get all documents for a course, optionally filtered by folder_id"""
        query = CSCLCourseDocument.query.filter_by(course_id=course_id)
        if folder_id:
            query = query.filter_by(folder_id=folder_id)
        documents = query.order_by(CSCLCourseDocument.created_at.desc()).all()
        
        result = []
        for doc in documents:
            doc_dict = doc.to_dict(include_file_size=True)
            doc_dict['filename'] = doc.title
            doc_dict['chunks_count'] = len(doc.chunks)
            # Extract preview from chunks (only extracted text, never raw bytes)
            # S2.13: skip any chunk that looks like PDF binary (e.g. legacy bad data)
            if doc.chunks:
                safe_chunks = [c.chunk_text for c in doc.chunks if c.chunk_text and not is_probably_pdf_binary_text(c.chunk_text)]
                full_text = ' '.join(safe_chunks) if safe_chunks else ''
                # Nuclear filter: only send printable-safe preview to client (no binary can slip through)
                preview = to_display_safe_preview(full_text, _PREVIEW_LEN)
                doc_dict['extracted_text_preview'] = preview  # None -> frontend shows "未提取到文本"
            else:
                doc_dict['extracted_text_preview'] = None
            result.append(doc_dict)
        
        return result
    
    def delete_document(self, document_id: str, course_id: str, user_id: str) -> bool:
        """Delete a document (only if user uploaded it or is admin)"""
        document = CSCLCourseDocument.query.filter_by(
            id=document_id,
            course_id=course_id
        ).first()
        
        if not document:
            return False
        
        # Check permission (uploader or admin)
        from app.models import User, UserRole
        user = User.query.get(user_id)
        if document.uploaded_by != user_id and (not user or user.role != UserRole.ADMIN):
            return False
        
        self._cleanup_temp_file(document.storage_uri)
        
        # Delete document (chunks will be cascade deleted)
        db.session.delete(document)
        db.session.commit()
        
        return True
